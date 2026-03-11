"""Document processing service for text extraction."""

import gc
import logging
import threading
from pathlib import Path
from typing import Callable

logger = logging.getLogger(__name__)


class ProcessingCancelledError(Exception):
    """Raised when document processing is cancelled."""
    pass


# Singleton for OCR model to avoid reloading
_ocr_model = None
_ocr_processor = None
_ocr_model_lock = threading.Lock()
_ocr_loaded_model_id = None  # Track which model is loaded


def _get_device():
    """Get the best available device for inference."""
    import torch
    if torch.cuda.is_available():
        return "cuda"
    elif torch.backends.mps.is_available():
        return "mps"
    return "cpu"


def _get_active_ocr_model_id() -> str | None:
    """Get the active OCR model ID from the catalog."""
    try:
        from core.ocr_catalog import _read_active_ocr_model
        return _read_active_ocr_model()
    except ImportError:
        return None


def _get_ocr_model_path() -> str | None:
    """Get the path to the active downloaded OCR model."""
    try:
        from core.ocr_catalog import get_model_path, is_model_downloaded, _read_active_ocr_model
        active_id = _read_active_ocr_model()
        if active_id and is_model_downloaded(active_id):
            path = get_model_path(active_id)
            return str(path) if path else None
    except ImportError:
        pass
    return None


def _get_ocr_architecture() -> str | None:
    """Get the architecture of the active OCR model."""
    try:
        from core.ocr_catalog import OCR_MODEL_CATALOG, _read_active_ocr_model
        active_id = _read_active_ocr_model()
        if active_id:
            entry = OCR_MODEL_CATALOG.get(active_id)
            if entry:
                return entry.get("architecture")
    except ImportError:
        pass
    return None


def _is_ocr_model_ready() -> bool:
    """Check if an OCR model is downloaded and ready."""
    try:
        from core.ocr_catalog import is_model_downloaded, _read_active_ocr_model
        active_id = _read_active_ocr_model()
        if active_id:
            return is_model_downloaded(active_id)
    except ImportError:
        pass
    return False


def _load_qwen2_vl(model_path: str, device: str):
    """Load Qwen2-VL model and processor."""
    import torch
    from transformers import Qwen2VLForConditionalGeneration, AutoProcessor

    processor = AutoProcessor.from_pretrained(
        model_path,
        trust_remote_code=True,
    )

    if device == "cuda":
        model = Qwen2VLForConditionalGeneration.from_pretrained(
            model_path,
            torch_dtype=torch.bfloat16,
            device_map="auto",
            trust_remote_code=True,
        )
    elif device == "mps":
        model = Qwen2VLForConditionalGeneration.from_pretrained(
            model_path,
            torch_dtype=torch.float16,
            trust_remote_code=True,
        ).to(device)
    else:
        model = Qwen2VLForConditionalGeneration.from_pretrained(
            model_path,
            torch_dtype=torch.float32,
            trust_remote_code=True,
        )

    model.eval()
    return model, processor


def _load_granite_vision(model_path: str, device: str):
    """Load Granite Vision model and processor."""
    import torch
    from transformers import AutoModelForVision2Seq, AutoProcessor

    processor = AutoProcessor.from_pretrained(
        model_path,
        trust_remote_code=True,
    )

    if device == "cuda":
        model = AutoModelForVision2Seq.from_pretrained(
            model_path,
            torch_dtype=torch.bfloat16,
            device_map="auto",
            trust_remote_code=True,
        )
    elif device == "mps":
        model = AutoModelForVision2Seq.from_pretrained(
            model_path,
            torch_dtype=torch.float16,
            trust_remote_code=True,
        ).to(device)
    else:
        model = AutoModelForVision2Seq.from_pretrained(
            model_path,
            torch_dtype=torch.float32,
            trust_remote_code=True,
        )

    model.eval()
    return model, processor


def _load_llama_vision(model_path: str, device: str):
    """Load Llama Vision model and processor."""
    import torch
    from transformers import MllamaForConditionalGeneration, AutoProcessor

    processor = AutoProcessor.from_pretrained(
        model_path,
        trust_remote_code=True,
    )

    if device == "cuda":
        model = MllamaForConditionalGeneration.from_pretrained(
            model_path,
            torch_dtype=torch.bfloat16,
            device_map="auto",
            trust_remote_code=True,
        )
    elif device == "mps":
        model = MllamaForConditionalGeneration.from_pretrained(
            model_path,
            torch_dtype=torch.float16,
            trust_remote_code=True,
        ).to(device)
    else:
        model = MllamaForConditionalGeneration.from_pretrained(
            model_path,
            torch_dtype=torch.float32,
            trust_remote_code=True,
        )

    model.eval()
    return model, processor


def get_ocr_model():
    """Get or create the singleton OCR model and processor."""
    global _ocr_model, _ocr_processor, _ocr_loaded_model_id
    with _ocr_model_lock:
        active_id = _get_active_ocr_model_id()

        # If the active model changed, clear the old one
        if _ocr_model is not None and _ocr_loaded_model_id != active_id:
            logger.info("Active OCR model changed from %s to %s, reloading...", _ocr_loaded_model_id, active_id)
            del _ocr_model
            del _ocr_processor
            _ocr_model = None
            _ocr_processor = None
            _ocr_loaded_model_id = None
            gc.collect()

        if _ocr_model is None:
            model_path = _get_ocr_model_path()
            if not model_path:
                raise RuntimeError("OCR model not downloaded. Download it in Settings > AI.")

            architecture = _get_ocr_architecture()
            device = _get_device()
            logger.info("Loading OCR model (arch=%s) from %s on %s...", architecture, model_path, device)

            if architecture == "qwen2-vl":
                _ocr_model, _ocr_processor = _load_qwen2_vl(model_path, device)
            elif architecture == "granite-vision":
                _ocr_model, _ocr_processor = _load_granite_vision(model_path, device)
            elif architecture == "llama-vision":
                _ocr_model, _ocr_processor = _load_llama_vision(model_path, device)
            else:
                # Fallback: try AutoModelForVision2Seq
                logger.warning("Unknown architecture '%s', attempting generic load", architecture)
                _ocr_model, _ocr_processor = _load_granite_vision(model_path, device)

            _ocr_loaded_model_id = active_id
            logger.info("OCR model loaded successfully (model_id=%s)", active_id)

        return _ocr_model, _ocr_processor


def cleanup_ocr_model():
    """Unload the OCR model to free memory."""
    global _ocr_model, _ocr_processor, _ocr_loaded_model_id
    with _ocr_model_lock:
        if _ocr_model is not None:
            logger.info("Unloading OCR model to free memory...")

            del _ocr_model
            del _ocr_processor
            _ocr_model = None
            _ocr_processor = None
            _ocr_loaded_model_id = None

            # Multiple gc.collect() calls - some objects need multiple passes
            gc.collect()
            gc.collect()
            gc.collect()

            # Clear torch caches
            try:
                import torch
                if torch.backends.mps.is_available():
                    torch.mps.synchronize()
                    torch.mps.empty_cache()
                elif torch.cuda.is_available():
                    torch.cuda.synchronize()
                    torch.cuda.empty_cache()
            except Exception as e:
                logger.debug(f"Error clearing torch cache: {e}")

            gc.collect()
            logger.info("OCR model unloaded")


def _format_ocr_messages_qwen2vl(image):
    """Format OCR prompt for Qwen2-VL models."""
    return [
        {
            "role": "user",
            "content": [
                {"type": "image", "image": image},
                {"type": "text", "text": "Extract and transcribe all text from this image. Preserve the layout and formatting as much as possible."},
            ],
        }
    ]


def _run_ocr_qwen2vl(image, model, processor, device, check_cancelled):
    """Run OCR using Qwen2-VL model."""
    import torch
    from qwen_vl_utils import process_vision_info

    messages = _format_ocr_messages_qwen2vl(image)

    text = processor.apply_chat_template(
        messages, tokenize=False, add_generation_prompt=True
    )
    image_inputs, video_inputs = process_vision_info(messages)

    inputs = processor(
        text=[text],
        images=image_inputs,
        videos=video_inputs,
        padding=True,
        return_tensors="pt",
    )

    if device != "cpu":
        inputs = inputs.to(device)

    if check_cancelled and check_cancelled():
        raise ProcessingCancelledError("Processing cancelled during preparation")

    with torch.no_grad():
        generated_ids = model.generate(**inputs, max_new_tokens=2048, do_sample=False)

    generated_ids_trimmed = [
        out_ids[len(in_ids):] for in_ids, out_ids in zip(inputs.input_ids, generated_ids)
    ]
    output_text = processor.batch_decode(
        generated_ids_trimmed, skip_special_tokens=True, clean_up_tokenization_spaces=False
    )[0]

    del inputs, generated_ids, generated_ids_trimmed, image_inputs, video_inputs
    return output_text


def _run_ocr_generic(image, model, processor, device, check_cancelled):
    """Run OCR using Granite Vision or Llama Vision models (standard HF transformers API)."""
    import torch

    prompt = "Extract and transcribe all text from this image. Preserve the layout and formatting as much as possible."

    messages = [
        {
            "role": "user",
            "content": [
                {"type": "image"},
                {"type": "text", "text": prompt},
            ],
        }
    ]

    input_text = processor.apply_chat_template(
        messages, add_generation_prompt=True, tokenize=False
    )

    inputs = processor(
        text=input_text,
        images=[image],
        return_tensors="pt",
    )

    if device != "cpu":
        inputs = inputs.to(device)

    if check_cancelled and check_cancelled():
        raise ProcessingCancelledError("Processing cancelled during preparation")

    with torch.no_grad():
        generated_ids = model.generate(**inputs, max_new_tokens=2048, do_sample=False)

    # Trim input tokens from output
    generated_ids_trimmed = generated_ids[:, inputs["input_ids"].shape[-1]:]
    output_text = processor.batch_decode(
        generated_ids_trimmed, skip_special_tokens=True, clean_up_tokenization_spaces=False
    )[0]

    del inputs, generated_ids, generated_ids_trimmed
    return output_text


def _run_ocr_on_image(image, check_cancelled: Callable[[], bool] | None = None) -> str:
    """Run OCR on a single image using the active vision model."""
    import torch

    if check_cancelled and check_cancelled():
        raise ProcessingCancelledError("Processing cancelled before OCR")

    model, processor = get_ocr_model()
    device = _get_device()
    architecture = _get_ocr_architecture()

    if architecture == "qwen2-vl":
        output_text = _run_ocr_qwen2vl(image, model, processor, device, check_cancelled)
    else:
        output_text = _run_ocr_generic(image, model, processor, device, check_cancelled)

    # Clear device cache
    if device == "mps":
        torch.mps.empty_cache()
    elif device == "cuda":
        torch.cuda.empty_cache()

    # Clean up any remaining special tokens that weren't filtered
    for token in ["<|im_end|>", "<|im_start|>", "<|endoftext|>", "<|eot_id|>"]:
        output_text = output_text.replace(token, "")

    return output_text.strip()


def _check_pymupdf_available() -> bool:
    """Check if PyMuPDF is installed."""
    try:
        import fitz
        return True
    except ImportError:
        return False


def _check_pypdf_available() -> bool:
    """Check if pypdf is installed."""
    try:
        from pypdf import PdfReader
        return True
    except ImportError:
        return False


PYMUPDF_AVAILABLE = _check_pymupdf_available()
PYPDF_AVAILABLE = _check_pypdf_available()


class DocumentProcessor:
    """Extracts text from various document formats."""

    def _is_ocr_available(self) -> bool:
        """Check if OCR model is downloaded and ready."""
        return _is_ocr_model_ready()

    def process(
        self,
        file_path: Path,
        mime_type: str,
        enable_ocr: bool = False,
        check_cancelled: Callable[[], bool] | None = None,
    ) -> dict:
        """
        Process a document and extract text content.

        Args:
            file_path: Path to the document file
            mime_type: MIME type of the document
            enable_ocr: If True, use OCR for text extraction (for scanned docs/images)
            check_cancelled: Optional callback that returns True if processing should stop

        Returns:
            dict with keys: text, markdown, page_count, metadata

        Raises:
            ProcessingCancelledError: If processing was cancelled
        """
        if mime_type == "application/pdf":
            return self._process_pdf(file_path, enable_ocr=enable_ocr, check_cancelled=check_cancelled)
        elif mime_type.startswith("image/"):
            return self._process_image(file_path, enable_ocr=enable_ocr, check_cancelled=check_cancelled)
        elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            return self._process_docx(file_path, enable_ocr=enable_ocr)
        elif mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            return self._process_xlsx(file_path)
        elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            return self._process_pptx(file_path, enable_ocr=enable_ocr)
        elif mime_type in ("text/plain", "text/markdown"):
            return self._process_text(file_path)
        else:
            raise ValueError(f"Unsupported MIME type: {mime_type}")

    def _process_pdf(
        self,
        file_path: Path,
        enable_ocr: bool = False,
        check_cancelled: Callable[[], bool] | None = None,
    ) -> dict:
        """Process PDF: try text extraction first, use OCR only for scanned/image PDFs."""
        if check_cancelled and check_cancelled():
            raise ProcessingCancelledError("Processing cancelled before starting")

        # Always try PyMuPDF text extraction first (fast and accurate for text-based PDFs)
        pymupdf_result = self._process_pdf_fallback(file_path)

        # Check if we got meaningful text (threshold: avg 100 chars per page)
        text_length = len(pymupdf_result.get("text", "") or "")
        page_count = pymupdf_result.get("page_count") or 1
        chars_per_page = text_length / page_count if page_count > 0 else 0

        # If we got good text extraction, use it (no need for slow OCR)
        if chars_per_page >= 100:
            logger.info(f"PyMuPDF extracted {text_length} chars ({chars_per_page:.0f}/page) from {file_path.name}")
            return pymupdf_result

        # Text extraction yielded sparse results - this might be a scanned PDF
        logger.info(f"Sparse text extraction ({chars_per_page:.0f} chars/page) for {file_path.name}")

        # Only use OCR if explicitly enabled AND model is available
        if not enable_ocr:
            logger.info(f"OCR not enabled, returning sparse PyMuPDF result for {file_path.name}")
            return pymupdf_result

        if not self._is_ocr_available():
            logger.info(f"OCR enabled but model not downloaded for {file_path.name}")
            return pymupdf_result

        # Run OCR on the scanned/image PDF
        try:
            logger.info(f"Running OCR on scanned PDF: {file_path.name}")

            if not PYMUPDF_AVAILABLE:
                logger.warning("PyMuPDF not available for PDF to image conversion")
                return pymupdf_result

            import fitz
            from PIL import Image
            import io

            doc = fitz.open(file_path)
            markdown_parts = []

            for i, page in enumerate(doc):
                if check_cancelled and check_cancelled():
                    doc.close()
                    raise ProcessingCancelledError(f"Processing cancelled at page {i+1}/{len(doc)}")

                logger.info(f"OCR processing page {i+1}/{len(doc)} of {file_path.name}")

                # Render page to image (150 DPI - balance between quality and memory)
                mat = fitz.Matrix(150/72, 150/72)
                pix = page.get_pixmap(matrix=mat)
                img_data = pix.tobytes("png")

                # Clean up pixmap immediately to free memory
                del pix

                image = Image.open(io.BytesIO(img_data))
                del img_data  # Free the PNG bytes

                # Run OCR on the page image
                page_text = _run_ocr_on_image(image, check_cancelled)
                markdown_parts.append(f"## Page {i+1}\n\n{page_text}")

                # Clean up image and force garbage collection
                del image
                gc.collect()

            doc.close()

            active_id = _get_active_ocr_model_id() or "ocr"
            combined_markdown = "\n\n".join(markdown_parts)
            plain_text = combined_markdown.replace("#", "").replace("*", "").replace("|", " ")

            return {
                "text": plain_text,
                "markdown": combined_markdown,
                "page_count": len(markdown_parts),
                "metadata": {"ocr_engine": active_id},
            }
        except ProcessingCancelledError:
            raise
        except Exception as e:
            logger.warning(f"OCR failed, using PyMuPDF result: {e}")
            return pymupdf_result

    def _process_pdf_fallback(self, file_path: Path) -> dict:
        """Fallback PDF processing using PyMuPDF or pypdf."""
        if PYMUPDF_AVAILABLE:
            try:
                import fitz
                logger.info(f"Processing {file_path.name} with PyMuPDF")
                doc = fitz.open(file_path)
                text_parts = []
                markdown_parts = []
                for i, page in enumerate(doc):
                    page_text = page.get_text()
                    text_parts.append(page_text)
                    markdown_parts.append(f"## Page {i+1}\n\n{page_text}")
                text = "\n\n".join(text_parts)
                markdown = "\n\n".join(markdown_parts)
                page_count = len(doc)
                doc.close()
                return {
                    "text": text,
                    "markdown": markdown,
                    "page_count": page_count,
                    "metadata": {"ocr_engine": "pymupdf"},
                }
            except Exception as e:
                logger.warning(f"PyMuPDF processing failed, trying pypdf: {e}")

        if PYPDF_AVAILABLE:
            try:
                from pypdf import PdfReader
                logger.info(f"Processing {file_path.name} with pypdf")
                reader = PdfReader(file_path)
                text_parts = []
                markdown_parts = []
                for i, page in enumerate(reader.pages):
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(page_text)
                        markdown_parts.append(f"## Page {i+1}\n\n{page_text}")
                text = "\n\n".join(text_parts)
                markdown = "\n\n".join(markdown_parts)
                return {
                    "text": text,
                    "markdown": markdown,
                    "page_count": len(reader.pages),
                    "metadata": {"ocr_engine": "pypdf"},
                }
            except Exception as e:
                logger.error(f"pypdf processing failed: {e}")

        logger.error("No PDF processing library available")
        return {
            "text": "",
            "markdown": "",
            "page_count": None,
            "metadata": {"error": "No PDF processor available"},
        }

    def _process_image(
        self,
        file_path: Path,
        enable_ocr: bool = False,
        check_cancelled: Callable[[], bool] | None = None,
    ) -> dict:
        """Process image using OCR if enabled."""
        if not enable_ocr:
            logger.debug(f"OCR not enabled for image: {file_path.name}")
            return {
                "text": "",
                "markdown": "",
                "page_count": 1,
                "metadata": {"format": "image", "ocr_enabled": False},
            }

        if not self._is_ocr_available():
            logger.warning(f"OCR model not downloaded for image: {file_path.name}")
            return {
                "text": "",
                "markdown": "",
                "page_count": 1,
                "metadata": {"error": "OCR model not downloaded. Download it in Settings > AI.", "ocr_enabled": True},
            }

        try:
            if check_cancelled and check_cancelled():
                raise ProcessingCancelledError("Processing cancelled before starting")

            from PIL import Image

            active_id = _get_active_ocr_model_id() or "ocr"
            logger.info(f"Processing {file_path.name} with OCR model {active_id}")

            # Load image
            image = Image.open(file_path)

            # Convert to RGB if necessary (handles RGBA, grayscale, etc.)
            if image.mode != "RGB":
                image = image.convert("RGB")

            # Resize large images to prevent OOM (phone photos can be 4000+ px)
            max_dim = 1280
            if max(image.size) > max_dim:
                image.thumbnail((max_dim, max_dim), Image.LANCZOS)
                logger.info(f"Resized image to {image.size[0]}x{image.size[1]} for OCR")

            # Run OCR
            text = _run_ocr_on_image(image, check_cancelled)

            return {
                "text": text,
                "markdown": text,
                "page_count": 1,
                "metadata": {"ocr_engine": active_id},
            }
        except ProcessingCancelledError:
            raise
        except Exception as e:
            logger.error(f"Image OCR failed: {e}")
            return {"text": "", "markdown": "", "page_count": 1, "metadata": {"error": str(e)}}

    def _process_docx(self, file_path: Path, enable_ocr: bool = False) -> dict:
        """Process DOCX using python-docx."""
        try:
            from docx import Document as DocxDocument
            from docx.oxml.ns import qn

            doc = DocxDocument(file_path)
            markdown_parts = []
            text_parts = []

            def extract_textbox_text(para) -> str:
                """Extract text from textboxes/shapes in a paragraph."""
                texts = []
                seen = set()
                for run in para.runs:
                    for txbx in run._element.findall('.//' + qn('w:txbxContent')):
                        txbx_text = []
                        for t_elem in txbx.findall('.//' + qn('w:t')):
                            if t_elem.text:
                                txbx_text.append(t_elem.text)
                        full_text = ''.join(txbx_text).strip()
                        if full_text and full_text not in seen:
                            seen.add(full_text)
                            texts.append(full_text)
                return ' '.join(texts).strip()

            for element in doc.element.body:
                if element.tag.endswith('p'):
                    for para in doc.paragraphs:
                        if para._element is element:
                            textbox_text = extract_textbox_text(para)
                            if textbox_text:
                                markdown_parts.append(f"# {textbox_text}")
                                text_parts.append(textbox_text)

                            para_text = para.text.strip()
                            if para_text and para_text != textbox_text:
                                if para.style and para.style.name.startswith('Heading'):
                                    level = para.style.name[-1] if para.style.name[-1].isdigit() else '2'
                                    markdown_parts.append(f"{'#' * int(level)} {para_text}")
                                else:
                                    markdown_parts.append(para_text)
                                text_parts.append(para_text)
                            break
                elif element.tag.endswith('tbl'):
                    for table in doc.tables:
                        if table._element is element:
                            table_md = self._table_to_markdown(table)
                            if table_md:
                                markdown_parts.append(table_md)
                                for row in table.rows:
                                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                                    if row_text.replace("|", "").strip():
                                        text_parts.append(row_text)
                            break

            for section in doc.sections:
                footer = section.footer
                for para in footer.paragraphs:
                    if para.text.strip():
                        markdown_parts.append(f"*{para.text}*")
                        text_parts.append(para.text)

            return {
                "text": "\n\n".join(text_parts),
                "markdown": "\n\n".join(markdown_parts),
                "page_count": None,
                "metadata": {"format": "docx"},
            }
        except ImportError:
            logger.warning("python-docx not installed")
            return {"text": "", "markdown": "", "page_count": None, "metadata": {}}
        except Exception as e:
            logger.error(f"DOCX processing failed: {e}")
            return {"text": "", "markdown": "", "page_count": None, "metadata": {}}

    def _table_to_markdown(self, table) -> str:
        """Convert a DOCX table to markdown format."""
        rows = []
        for i, row in enumerate(table.rows):
            cells = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
            if not any(cells):
                continue
            rows.append("| " + " | ".join(cells) + " |")
            if i == 0:
                rows.append("| " + " | ".join("---" for _ in cells) + " |")
        return "\n".join(rows) if rows else ""

    def _process_xlsx(self, file_path: Path) -> dict:
        """Process XLSX using openpyxl."""
        try:
            from openpyxl import load_workbook
            wb = load_workbook(file_path, read_only=True, data_only=True)
            markdown_parts = []
            text_parts = []

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                markdown_parts.append(f"## {sheet_name}\n")
                rows = list(sheet.iter_rows(values_only=True))
                if rows:
                    header = rows[0]
                    markdown_parts.append("| " + " | ".join(str(c or "") for c in header) + " |")
                    markdown_parts.append("| " + " | ".join("---" for _ in header) + " |")
                    for row in rows[1:]:
                        markdown_parts.append("| " + " | ".join(str(c or "") for c in row) + " |")
                        text_parts.append("\t".join(str(c or "") for c in row))
                markdown_parts.append("")

            return {
                "text": "\n".join(text_parts),
                "markdown": "\n".join(markdown_parts),
                "page_count": len(wb.sheetnames),
                "metadata": {"format": "xlsx", "sheets": wb.sheetnames},
            }
        except ImportError:
            logger.warning("openpyxl not installed")
            return {"text": "", "markdown": "", "page_count": None, "metadata": {}}
        except Exception as e:
            logger.error(f"XLSX processing failed: {e}")
            return {"text": "", "markdown": "", "page_count": None, "metadata": {}}

    def _process_pptx(self, file_path: Path, enable_ocr: bool = False) -> dict:
        """Process PPTX using python-pptx."""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            markdown_parts = []
            text_parts = []

            for i, slide in enumerate(prs.slides, 1):
                markdown_parts.append(f"## Slide {i}\n")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        markdown_parts.append(shape.text)
                        text_parts.append(shape.text)
                markdown_parts.append("")

            return {
                "text": "\n\n".join(text_parts),
                "markdown": "\n".join(markdown_parts),
                "page_count": len(prs.slides),
                "metadata": {"format": "pptx"},
            }
        except ImportError:
            logger.warning("python-pptx not installed")
            return {"text": "", "markdown": "", "page_count": None, "metadata": {}}
        except Exception as e:
            logger.error(f"PPTX processing failed: {e}")
            return {"text": "", "markdown": "", "page_count": None, "metadata": {}}

    def _process_text(self, file_path: Path) -> dict:
        """Process plain text or markdown files."""
        try:
            text = file_path.read_text(encoding="utf-8")
            return {
                "text": text,
                "markdown": text,
                "page_count": 1,
                "metadata": {"format": "text"},
            }
        except Exception as e:
            logger.error(f"Text processing failed: {e}")
            return {"text": "", "markdown": "", "page_count": 1, "metadata": {}}


document_processor = DocumentProcessor()
